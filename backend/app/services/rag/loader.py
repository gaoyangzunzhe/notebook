"""文档加载器：按扩展名选择合适的加载器，产出 LangChain Document 列表。

支持格式：
- .pdf        PyPDFLoader
- .docx       python-docx（段落 + 表格）
- .pptx       python-pptx（逐页文本框 + 表格，带 [Slide N] 页码前缀）
- 其余        按纯文本读取（TextLoader）
"""
from pathlib import Path

from docx import Document as DocxDocument
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_core.documents import Document
from pptx import Presentation

from app.core.errors import RAGUnavailableError


def _load_docx(path: Path) -> list[Document]:
    """用 python-docx 提取 .docx 的段落与表格文本。"""
    doc = DocxDocument(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    content = "\n".join(parts)
    return [Document(page_content=content)] if content.strip() else []


def _load_pptx(path: Path) -> list[Document]:
    """用 python-pptx 逐页提取 .pptx 的文本框与表格文本。"""
    prs = Presentation(str(path))
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            slides_text.append(f"[Slide {i}]\n" + "\n\n".join(texts))
    content = "\n\n".join(slides_text)
    return [Document(page_content=content)] if content.strip() else []


def load_document(path: Path, filename: str) -> list[Document]:
    """读取单个文档，按扩展名分发到对应加载器。"""
    suffix = Path(filename).suffix.lower()
    try:
        if suffix == ".pdf":
            docs = PyPDFLoader(str(path)).load()
        elif suffix == ".docx":
            docs = _load_docx(path)
        elif suffix == ".pptx":
            docs = _load_pptx(path)
        else:
            docs = TextLoader(str(path), encoding="utf-8").load()
    except RAGUnavailableError:
        raise
    except Exception as e:  # noqa: BLE001
        raise RAGUnavailableError(f"无法读取文档 {filename}: {e}") from e

    # 覆盖 metadata.source：加载器默认记录存储路径，这里改成用户可读的原始文件名
    for doc in docs:
        doc.metadata["source"] = filename
    return docs
